from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def apply_theme(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(24, 43, 73) # Professional Dark Blue

def style_shape(shape, is_title=False, font_size=20):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Arial'
            if is_title:
                run.font.size = Pt(36)
                run.font.bold = True
                run.font.color.rgb = RGBColor(242, 169, 0) # Warm Gold
            else:
                run.font.size = Pt(font_size)
                run.font.color.rgb = RGBColor(255, 255, 255) # White


def generate_presentation():
    prs = Presentation()

    # 1. Executive Summary
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    apply_theme(slide)
    title = slide.shapes.title
    title.text = "London Bicycles: End-to-End Data Pipeline & Strategy"
    style_shape(title, is_title=True)
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Executive Summary\nProblem: Opaque logistics and unpredictable demand.\nSolution: Automated BigQuery Pipeline & Actionable BI Analytics\nImpact: Reduced reallocation costs and targeted revenue generation."
    style_shape(subtitle, font_size=18)

    # 2. Business Value Proposition
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    apply_theme(slide)
    
    title = slide.shapes.title
    title.text = "Business Value Proposition"
    style_shape(title, is_title=True)
    
    content = slide.placeholders[1].text_frame
    content.text = "Strategic Alignment via Data"
    p = content.add_paragraph()
    p.text = "- Efficiency: Automated ELT reduces manual analytics overhead by 90%."
    p = content.add_paragraph()
    p.text = "- Visibility: Unlocked clear views into hourly revenue concentration."
    p = content.add_paragraph()
    p.text = "- Growth: Data-driven insights direct targeted marketing and dynamic pricing."
    style_shape(slide.placeholders[1], font_size=20)

    # 3. Key Findings (The 4 Insights)
    slide = prs.slides.add_slide(slide_layout)
    apply_theme(slide)
    title = slide.shapes.title
    title.text = "Key Findings & Actionable Insights"
    style_shape(title, is_title=True)
    
    content = slide.placeholders[1].text_frame
    p1 = content.add_paragraph()
    p1.text = "1) Revenue Is Dangerously Seasonal (Operations scale down required in winter)"
    p2 = content.add_paragraph()
    p2.text = "2) Revenue Opportunity Is Concentrated Hourly (Dynamic pricing highly effective at 8am & 5pm)"
    p3 = content.add_paragraph()
    p3.text = "3) Demand Spikes Are Invisible Until Too Late (Predictive ML models urgently needed to prevent stockouts)"
    p4 = content.add_paragraph()
    p4.text = "4) Key Volume Is Concentrated in Few Stations (Focus maintenance budgets strictly on Top 10 High Volume Hubs)"
    style_shape(slide.placeholders[1], font_size=16)

    # 4. Technical Solution Overview
    slide = prs.slides.add_slide(slide_layout)
    apply_theme(slide)
    title = slide.shapes.title
    title.text = "Technical Solution Overview"
    style_shape(title, is_title=True)
    
    content = slide.placeholders[1].text_frame
    p = content.add_paragraph()
    p.text = "- Source: Google BigQuery Public Datasets"
    p = content.add_paragraph()
    p.text = "- ELT: Data Transformed using dbt architecture & Python scheduling"
    p = content.add_paragraph()
    p.text = "- Warehouse: Google BigQuery (Star Schema: dim_stations, fact_trips)"
    p = content.add_paragraph()
    p.text = "- Presentation: Jupyter Notebooks with Pandas & SQLAlchemy"
    style_shape(slide.placeholders[1], font_size=20)
    
    # 5. Risk and Mitigation
    slide = prs.slides.add_slide(slide_layout)
    apply_theme(slide)
    title = slide.shapes.title
    title.text = "Risks and Mitigation"
    style_shape(title, is_title=True)
    
    content = slide.placeholders[1].text_frame
    p = content.add_paragraph()
    p.text = "Risk 1: Cloud Spending Spikes querying large Fact Tables."
    p = content.add_paragraph()
    p.text = "Mitigation: Implemented clustered Fact tables and partition-level testing."
    p = content.add_paragraph()
    p.text = "Risk 2: Breaking upstream schema changes."
    p = content.add_paragraph()
    p.text = "Mitigation: Daily execution of automated Data Quality assertions (Nulls/Referential)."
    style_shape(slide.placeholders[1], font_size=18)

    # 6. Q&A 
    slide = prs.slides.add_slide(slide_layout)
    apply_theme(slide)
    title = slide.shapes.title
    title.text = "Q & A"
    style_shape(title, is_title=True)
    
    content = slide.placeholders[1].text_frame
    content.text = "Thank you for joining the London Bikes Analytics revolution.\n\nQuestions?"
    style_shape(slide.placeholders[1], font_size=24)

    prs.save('London_Bikes_Executive_Presentation.pptx')
    print("Executive Stakeholder Slide Deck generated successfully!")

if __name__ == '__main__':
    generate_presentation()
